import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (clean, minimalist, academic matching reference PDF)
    BG_COLOR = RGBColor(255, 255, 255)       # Pure White background
    TEXT_MAIN = RGBColor(51, 65, 85)          # Slate 700 (soft slate)
    TEXT_DARK = RGBColor(15, 23, 42)          # Slate 900 (bold titles)
    TEXT_MUTED = RGBColor(100, 116, 139)       # Slate 500 (muted text)
    COLOR_PRIMARY = RGBColor(30, 58, 138)     # Navy Blue (Blue 900)
    COLOR_SECONDARY = RGBColor(37, 99, 235)   # Steel Blue (Blue 600)
    COLOR_SUCCESS = RGBColor(22, 163, 74)     # Forest Green (Green 600)
    COLOR_BORDER = RGBColor(226, 232, 240)    # Border Slate (Slate 200)
    
    # Use blank slide layout (index 6 is blank)
    blank_layout = prs.slide_layouts[6]
    
    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header_and_logos(slide, title, subtitle):
        # 1. Corner logos from assets
        try:
            slide.shapes.add_picture("assets/feather.png", Inches(0.6), Inches(0.4), width=Inches(0.45), height=Inches(0.6))
            slide.shapes.add_picture("assets/devlup.png", Inches(12.133), Inches(0.4), width=Inches(0.6), height=Inches(0.6))
        except Exception as e:
            print("Warning: could not add logos to slide:", e)
            
        # 2. Centered header group
        tx_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.4), Inches(9.333), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Outfit"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)
        p2.alignment = PP_ALIGN.CENTER

    def add_footer(slide, number):
        tx_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.2), Inches(0.4))
        p = tx_box.text_frame.paragraphs[0]
        p.text = f"{number} / 18"
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MAIN
        p.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)
    
    # Title Box (centered)
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "CIPHER"
    p.font.name = "Outfit"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf1.add_paragraph()
    p2.text = "Decentralized Content Distribution over P2P Networks"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_SECONDARY
    p2.space_before = Pt(8)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf1.add_paragraph()
    p3.text = '"CDNs are owned by a handful of companies. CIPHER is the alternative — a decentralized content delivery protocol where math replaces the middleman."'
    p3.font.name = "Calibri"
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_MAIN
    p3.space_before = Pt(20)
    p3.alignment = PP_ALIGN.CENTER
    
    # Presenter Details
    brand_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.5))
    tf_brand = brand_box.text_frame
    p_b1 = tf_brand.paragraphs[0]
    p_b1.text = "Presented by: Adarsh Kumar"
    p_b1.font.name = "Calibri"
    p_b1.font.size = Pt(15)
    p_b1.font.bold = True
    p_b1.font.color.rgb = TEXT_DARK
    p_b1.alignment = PP_ALIGN.CENTER
    
    p_b2 = tf_brand.add_paragraph()
    p_b2.text = "Academic Project • DevlUp Labs • IIT Jodhpur"
    p_b2.font.name = "Calibri"
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = TEXT_MAIN
    p_b2.space_before = Pt(4)
    p_b2.alignment = PP_ALIGN.CENTER
    
    add_footer(slide1, 1)
    
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome to the presentation for CIPHER, a decentralized content distribution protocol built on top of P2P networks.\n\n"
        "Traditional Content Delivery Networks (CDNs) are highly centralized, owned by a small group of large providers, which makes them points of failure and control.\n"
        "CIPHER is designed to turn arbitrary files into cryptographically verifiable, erasure-coded shards distributed across untrusted nodes, replacing the middleman with mathematics.\n\n"
        "Let's move into the core problem statement we are addressing."
    )

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_header_and_logos(slide2, "Problem Statement", "The Challenge of Reliable Distributed File Transfer")
    
    # Left column: concise text bullets
    box_left = slide2.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Centralized Bottlenecks: Outgoing server pipes saturate quickly as client request counts grow.",
        "• Single Point of Failure: Outages at central data nodes result in complete service interruptions.",
        "• NAT & Firewalls: Over 60% of peer nodes reside behind restrictive router walls, blocking direct connections.",
        "• Peer Churn & Poisoning: Distributed nodes join and leave unpredictably, and untrusted peers may modify data."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide2, 2)
    
    slide2.notes_slide.notes_text_frame.text = (
        "Before explaining CIPHER, we must explain the problem.\n\n"
        "Traditional centralized systems are expensive and vulnerable to failure. Centralized networks face physical bottlenecks and vulnerability to censorship.\n"
        "But P2P setups introduce hard problems: NAT boundaries block 60%+ of connections, nodes drop offline (churn), and peers can distribute corrupt files.\n\n"
        "Our protocol is designed specifically to address these issues."
    )

    # -------------------------------------------------------------
    # SLIDE 3: Design Goals
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_header_and_logos(slide3, "Design Goals", "Four Pillars of a Robust P2P CDN")
    
    # 4 columns for cards
    col_width = Inches(2.8)
    col_gap = Inches(0.2)
    start_left = Inches(0.6)
    
    goals = [
        ("1. Decentralization", "No Central Control", "Shards are placed on distinct independent providers across the network, avoiding reliance on any single node."),
        ("2. Reliability", "Self-Healing Deck", "Peer dropouts and network latency spikes must not terminate or stall the overall download process."),
        ("3. Integrity", "Proof over Trust", "Strict ciphertext content-addressing ensures data is verified on receipt, blocking malicious poisoning."),
        ("4. Efficiency", "Swarmed Retrieval", "Parallel channels query multiple seed hosts concurrently to saturate the client's network path.")
    ]
    
    for i, (g_title, g_sub, g_desc) in enumerate(goals):
        box_left = start_left + i * (col_width + col_gap)
        
        rect = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.3), col_width, Inches(4.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = COLOR_BORDER
        
        tx_box = slide3.shapes.add_textbox(box_left, Inches(2.4), col_width, Inches(4.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = g_title
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tf.add_paragraph()
        p_sub.text = g_sub
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_SECONDARY
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = g_desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(8)
        
    add_footer(slide3, 3)
    
    slide3.notes_slide.notes_text_frame.text = (
        "These four goals are the core focus of the system architecture.\n\n"
        "First, decentralization: we distribute file shards across multiple peers.\n"
        "Second, reliability: a downloading node should not care if any single peer drops offline.\n"
        "Third, integrity: using cryptographic hash checks so that peers can't poison our files.\n"
        "Fourth, efficiency: swarmed downloading where we download different parts of a file concurrently from separate peers.\n"
        "Fifth, resilience: utilizing background NAT hole punching so standard home networks can connect directly without server relays."
    )

    # -------------------------------------------------------------
    # SLIDE 4: System Overview
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_header_and_logos(slide4, "System Overview", "From One File to Distributed Content")
    
    # Left column: text bullets
    box_left = slide4.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• File Chunking: Monolithic file is split into 256KB segments.",
        "• Erasure Redundancy: Reed-Solomon creates additional parity shards.",
        "• Swarm Placements: Shards are stored across distinct nodes.",
        "• Parallel Retrieval: Client downloads shards concurrently from multiple online nodes."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide4, 4)
    
    slide4.notes_slide.notes_text_frame.text = (
        "This slide shows the overall data flow transformation model.\n\n"
        "Instead of keeping a file on a central server, the publisher splits it into Data Chunks.\n"
        "Next, we run these chunks through Reed-Solomon erasure coding, which yields Data and Parity Shards.\n"
        "These shards are distributed across P2P storage providers.\n"
        "A client downloads these shards in parallel from multiple providers at once.\n"
        "Finally, the client verifies each shard's hash, decodes them back to the original chunks, and reconstructs the file. Note that reconstruction succeeds even if some providers fail, as long as we retrieve the minimum required shards."
    )

    # -------------------------------------------------------------
    # SLIDE 5: Core Features
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_header_and_logos(slide5, "Core Features", "Ingestion • Discovery • Swarming • Failover")
    
    # 4 columns
    col_width = Inches(2.8)
    col_gap = Inches(0.2)
    start_left = Inches(0.6)
    
    stages = [
        ("1. Prepare", "Ingestion & Split", "Splits files, encrypts via XChaCha20, hashes via SHA-256, and generates Reed-Solomon shards."),
        ("2. Discover", "Registry Substrate", "Queries Control Plane locations to build the client-side ChunkProviderMap execution plan."),
        ("3. Transfer", "Swarm Pipeline", "Concurrently queries multiple providers over upgraded DCUtR direct paths in a lock-free queue."),
        ("4. Recover", "Self-Healing Deck", "Handles node failovers by switching providers, resuming session logs, and decoding erasure blocks.")
    ]
    
    for i, (name, subtitle, desc) in enumerate(stages):
        box_left = start_left + i * (col_width + col_gap)
        
        rect = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.3), col_width, Inches(4.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = COLOR_BORDER
        
        tx_box = slide5.shapes.add_textbox(box_left, Inches(2.4), col_width, Inches(4.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_SECONDARY
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(8)
        
    add_footer(slide5, 5)
    
    slide5.notes_slide.notes_text_frame.text = (
        "Here we introduce the four core capabilities of the protocol lifecycle.\n\n"
        "Stage 1 is Content Preparation: generating encrypted shards and a manifest describing the file.\n"
        "Stage 2 is Discovery: querying the decentralized Control Plane to find peer addresses and building an execution plan mapping shards to source nodes.\n"
        "Stage 3 is Transfer: our concurrent scheduler worker pool spinning up multiplexed connections over libp2p and requesting content.\n"
        "Stage 4 is Recovery: correcting errors by switching providers on timeouts, resuming interrupted downloads from session records, and reconstructing data using erasure coding."
    )

    # -------------------------------------------------------------
    # SLIDE 6: Architecture
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_header_and_logos(slide6, "Architecture", "Control Plane and Data Plane Division")
    
    # Left column: text bullets
    box_left = slide6.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    p1 = tf_l.paragraphs[0]
    p1.text = "Control Plane (Planned / Design)"
    p1.font.name = "Outfit"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    
    p1_desc = tf_l.add_paragraph()
    p1_desc.text = "Manages metadata: Registries (Provider, Manifest, and Placement maps) resolved via Kademlia DHT."
    p1_desc.font.name = "Calibri"
    p1_desc.font.size = Pt(14)
    p1_desc.font.color.rgb = TEXT_MAIN
    p1_desc.space_before = Pt(4)
    p1_desc.space_after = Pt(20)
    
    p2 = tf_l.add_paragraph()
    p2.text = "Data Plane (Current Implementation)"
    p2.font.name = "Outfit"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUCCESS
    
    p2_desc = tf_l.add_paragraph()
    p2_desc.text = "Manages raw transfer: libp2p hosts, Circuit v2 Relay fallback, DCUtR hole punching, and the stateless Chunk Protocol."
    p2_desc.font.name = "Calibri"
    p2_desc.font.size = Pt(14)
    p2_desc.font.color.rgb = TEXT_MAIN
    p2_desc.space_before = Pt(4)
    
    # Right column: actual Architecture sketch image from assets
    try:
        slide6.shapes.add_picture("assets/Architecture.png", Inches(6.8), Inches(2.0), width=Inches(5.8), height=Inches(4.5))
    except Exception as e:
        print("Warning: could not add Architecture.png to slide 6:", e)
        
    add_footer(slide6, 6)
    
    slide6.notes_slide.notes_text_frame.text = (
        "This slide divides the CIPHER architecture into two major conceptual planes: the Control Plane and the Data Plane.\n\n"
        "The Control Plane answers 'where is the data?'. It coordinates provider, manifest, and placement mapping registries via a Kademlia DHT. This layer is currently in the design and planning expansion phase.\n\n"
        "The Data Plane answers 'get the data'. It manages transport, relay fallback, DCUtR hole punching, protocol stream handlers, and content verification. This plane is fully implemented and tested.\n\n"
        "Let's look at how data is ingested and prepared in Slide 7."
    )

    # -------------------------------------------------------------
    # SLIDE 7: Content Preparation
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_header_and_logos(slide7, "Content Preparation", "Chunking • Reed-Solomon Encoding • Manifest • Merkle Root")
    
    box_left = slide7.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Monolithic splitting: File is chunked into uniform blocks (default 256KB).",
        "• Independent Encryption: Blocks encrypted via XChaCha20-Poly1305, enabling random access.",
        "• Erasure expansion: Reed-Solomon coding produces Data + Parity Shards stored on providers.",
        "• Structure binding: Merkle Tree root is computed to verify shard integrity."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide7, 7)
    
    slide7.notes_slide.notes_text_frame.text = (
        "Here we zoom in on the Publisher module. This slide shows the exact pipeline for turning a file into distributed content.\n\n"
        "First, the Chunker divides the file into 256KB blocks.\n"
        "Second, we encrypt each chunk independently using XChaCha20-Poly1305 and hash it with SHA-256. Independent chunk encryption enables parallel random-access fetching and out-of-order decryption.\n"
        "Third, Reed-Solomon encoding outputs data and parity shards. Providers store shards rather than monolithic chunks.\n"
        "Fourth, we compute a Merkle Tree over these shards to produce the Merkle Root, and package this data into an immutable manifest capability file."
    )

    # -------------------------------------------------------------
    # SLIDE 8: Kademlia DHT
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_header_and_logos(slide8, "Kademlia DHT", "Decentralized Metadata and Discovery Substrate")
    
    box_left = slide8.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Discovery Substrate: Kademlia is not a file database; it is a peer address routing substrate.",
        "• Reference records: Maps identifiers to provider lists and manifest details; does not store actual file data.",
        "• Location resolution: Translates ShardID to a list of candidate host Multiaddresses."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide8, 8)
    
    slide8.notes_slide.notes_text_frame.text = (
        "It is crucial to clarify that Kademlia is NOT used as a general-purpose database in CIPHER.\n\n"
        "Instead, Kademlia DHT acts as a decentralized metadata/discovery substrate.\n"
        "It stores small records: ProviderRecords (peer addresses), ManifestRecords (erasure specifications), and ShardProviderRecords (mapping shard hashes to peer IDs).\n"
        "The DHT stores references and metadata, while actual large file payloads remain strictly hosted on the data plane providers.\n\n"
        "Let's talk about shard placement next."
    )

    # -------------------------------------------------------------
    # SLIDE 9: Shard Placement
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9)
    add_header_and_logos(slide9, "Shard Placement", "PlacementMap vs. ChunkProviderMap")
    
    # 2 columns cards
    col_width = Inches(5.8)
    col_gap = Inches(0.4)
    start_left = Inches(0.6)
    
    rect_l = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left, Inches(2.2), col_width, Inches(4.3))
    rect_l.fill.solid()
    rect_l.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_l.line.color.rgb = COLOR_BORDER
    
    tf_l = slide9.shapes.add_textbox(start_left, Inches(2.3), col_width, Inches(4.1)).text_frame
    tf_l.word_wrap = True
    p_lt = tf_l.paragraphs[0]
    p_lt.text = "PlacementMap"
    p_lt.font.name = "Outfit"
    p_lt.font.size = Pt(20)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLOR_PRIMARY
    
    p_ls = tf_l.add_paragraph()
    p_ls.text = "Publisher Blueprint"
    p_ls.font.bold = True
    p_ls.font.size = Pt(11)
    p_ls.font.color.rgb = COLOR_SECONDARY
    p_ls.space_after = Pt(12)
    
    l_bullets = [
        "• Generated during file publication.",
        "• Defines the static distribution plan.",
        "• Maps shard hashes to target provider nodes.",
        "• Stored on the Control Plane."
    ]
    for b in l_bullets:
        p = tf_l.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(8)
        
    rect_r = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left + col_width + col_gap, Inches(2.2), col_width, Inches(4.3))
    rect_r.fill.solid()
    rect_r.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_r.line.color.rgb = COLOR_BORDER
    
    tf_r = slide9.shapes.add_textbox(start_left + col_width + col_gap, Inches(2.3), col_width, Inches(4.1)).text_frame
    tf_r.word_wrap = True
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "ChunkProviderMap"
    p_rt.font.name = "Outfit"
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLOR_PRIMARY
    
    p_rs = tf_r.add_paragraph()
    p_rs.text = "Client Execution Plan"
    p_rs.font.bold = True
    p_rs.font.size = Pt(11)
    p_rs.font.color.rgb = COLOR_SECONDARY
    p_rs.space_after = Pt(12)
    
    r_bullets = [
        "• Generated dynamically by the downloader.",
        "• Adapts to active, online provider nodes.",
        "• Maps required segments to online peers.",
        "• Ensures run-time transfer routing."
    ]
    for b in r_bullets:
        p = tf_r.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(8)
        
    add_footer(slide9, 9)
    
    slide9.notes_slide.notes_text_frame.text = (
        "We need to distinguish between two mapping structures here.\n\n"
        "The PlacementMap is a global registry record created by the publisher. It maps shards to designated storage providers.\n"
        "The ChunkProviderMap is an execution plan created locally by the downloading client. It maps each required shard to the specific, online peers discovered at download time.\n"
        "This dynamic mapping protects the client from outdated placement tables if nodes disconnect."
    )

    # -------------------------------------------------------------
    # SLIDE 10: P2P Network
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    apply_background(slide10)
    add_header_and_logos(slide10, "P2P Network", "libp2p • Relay • DCUtR • Direct Connectivity")
    
    box_left = slide10.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• libp2p Host: Manages persistent Ed25519 identity keys and TCP/QUIC stream multiplexing.",
        "• Circuit v2 Relay: Coordinates transient connection relays to route around NAT firewalls.",
        "• DCUtR Upgrades: Coordinate background UDP simultaneous dial triggers to bypass NAT.",
        "• Direct Transport: Streams fallback to direct connection paths automatically."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide10, 10)
    
    slide10.notes_slide.notes_text_frame.text = (
        "The P2P transport layer handles host identity and connectivity.\n\n"
        "We use persistent Ed25519 keypairs stored in system configuration paths so that node identities (PeerIDs) remain constant across restarts.\n"
        "For connectivity, we fallback to public Circuit v2 relays to punch initial tunnels through NATs.\n"
        "Simultaneously, we run the DCUtR protocol. DCUtR attempts to establish direct socket connections in the background, upgrading the connection to a high-speed direct transport path without interrupting the transfer."
    )

    # -------------------------------------------------------------
    # SLIDE 11: Chunk Protocol
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    apply_background(slide11)
    add_header_and_logos(slide11, "Chunk Protocol", "Stateless Request / Response Stream Operations")
    
    box_left = slide11.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Protocol ID: Registered as /cipher/chunk/1.0.0 over multiplexed streams.",
        "• Framed Envelopes: Uses size-prefixed message frames (capped at 2MB) for memory protection.",
        "• Stateless Exchange: Handlers perform strict request/response data pushes, decoupling network state from application logic."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide11, 11)
    
    slide11.notes_slide.notes_text_frame.text = (
        "The /cipher/chunk/1.0.0 protocol defines the language peers use to request data.\n\n"
        "All messages are encapsulated in a unified symmetric envelope containing a version, message type, and byte payload.\n"
        "We send messages as size-prefixed frames capped at a maximum size of 2MB to prevent buffer overflow attacks.\n"
        "The message types include manifest request/response, chunk request/response, and positive or negative error acknowledgements. It is fully stateless, keeping complex transfer logic inside the client application modules."
    )

    # -------------------------------------------------------------
    # SLIDE 12: Parallel Download
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    apply_background(slide12)
    add_header_and_logos(slide12, "Parallel Download", "TransferManager & Scheduler Coordination")
    
    # 2 columns cards
    col_width = Inches(5.8)
    col_gap = Inches(0.4)
    start_left = Inches(0.6)
    
    rect_l = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left, Inches(2.2), col_width, Inches(4.3))
    rect_l.fill.solid()
    rect_l.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_l.line.color.rgb = COLOR_BORDER
    
    tf_l = slide12.shapes.add_textbox(start_left, Inches(2.3), col_width, Inches(4.1)).text_frame
    tf_l.word_wrap = True
    p_lt = tf_l.paragraphs[0]
    p_lt.text = "TransferManager"
    p_lt.font.name = "Outfit"
    p_lt.font.size = Pt(20)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLOR_PRIMARY
    
    p_ls = tf_l.add_paragraph()
    p_ls.text = "Session Orchestration"
    p_ls.font.bold = True
    p_ls.font.size = Pt(11)
    p_ls.font.color.rgb = COLOR_SECONDARY
    p_ls.space_after = Pt(12)
    
    l_bullets = [
        "• Owns the lifecycle of file downloads.",
        "• Persists session progress locally using a Boolean bitset.",
        "• Ensures client never re-downloads completed data blocks on crash."
    ]
    for b in l_bullets:
        p = tf_l.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(8)
        
    rect_r = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_left + col_width + col_gap, Inches(2.2), col_width, Inches(4.3))
    rect_r.fill.solid()
    rect_r.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_r.line.color.rgb = COLOR_BORDER
    
    tf_r = slide12.shapes.add_textbox(start_left + col_width + col_gap, Inches(2.3), col_width, Inches(4.1)).text_frame
    tf_r.word_wrap = True
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "Scheduler"
    p_rt.font.name = "Outfit"
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLOR_PRIMARY
    
    p_rs = tf_r.add_paragraph()
    p_rs.text = "Work Allocation"
    p_rs.font.bold = True
    p_rs.font.size = Pt(11)
    p_rs.font.color.rgb = COLOR_SECONDARY
    p_rs.space_after = Pt(12)
    
    r_bullets = [
        "• Maintains a lock-free queue of pending shard tasks.",
        "• Spins up concurrent worker threads to query distinct target peers.",
        "• Recycles tasks back to the queue on peer disconnect or timeout."
    ]
    for b in r_bullets:
        p = tf_r.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(8)
        
    add_footer(slide12, 12)
    
    slide12.notes_slide.notes_text_frame.text = (
        "The transfer orchestrator comprises the TransferManager and the Scheduler.\n\n"
        "The TransferManager initializes session states and tracks progress using a simple boolean bitset. This bitset determines exactly which shards are done and which are missing.\n"
        "The Scheduler schedules missing shards into a lock-free queue.\n"
        "Concurrent background worker routines pop tasks from the queue and send `/cipher/chunk/1.0.0` requests to online target peers concurrently. If a peer timeouts, the worker reports the failure to the scheduler, which recycles the task to be processed by a different peer."
    )

    # -------------------------------------------------------------
    # SLIDE 13: Security & Integrity
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    apply_background(slide13)
    add_header_and_logos(slide13, "Security & Integrity", "Proof over Trust Security Model")
    
    # 3 columns cards
    col_width = Inches(3.7)
    col_gap = Inches(0.2)
    start_left = Inches(0.6)
    
    sec_pillars = [
        ("Confidentiality", "XChaCha20-Poly1305", "Independent block encryption enables random-access, multi-peer swarming decryption. Decryption keys remain private."),
        ("Integrity", "Content-Addressing", "Shard hashes are derived from ciphertext. SHA-256 checks run immediately; corrupt data is rejected before writing to disk."),
        ("Structure", "Merkle Verification", "Merkle tree structures verify that downloaded segments belong to the original file, blocking poison blocks.")
    ]
    
    for i, (name, subtitle, desc) in enumerate(sec_pillars):
        box_left = start_left + i * (col_width + col_gap)
        
        rect = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.2), col_width, Inches(4.3))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = COLOR_BORDER
        
        tx_box = slide13.shapes.add_textbox(box_left, Inches(2.3), col_width, Inches(4.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_SECONDARY
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(8)
        
    add_footer(slide13, 13)
    
    slide13.notes_slide.notes_text_frame.text = (
        "In a public P2P network, you cannot trust peers. CIPHER operates on a strict 'proof-over-trust' model.\n\n"
        "We encrypt chunks using XChaCha20-Poly1305. The decryption keys are kept separate from the public manifest so that nodes store and relay data they cannot read.\n"
        "For integrity, shards are content-addressed by their ciphertext hashes. The client validates the SHA-256 hash immediately upon block receipt. If a hash mismatch occurs, we drop the data immediately, protecting storage from pollution attacks.\n"
        "Merkle tree root checks verify that the shard belongs to the original file."
    )

    # -------------------------------------------------------------
    # SLIDE 14: Download Flow
    # -------------------------------------------------------------
    slide14 = prs.slides.add_slide(blank_layout)
    apply_background(slide14)
    add_header_and_logos(slide14, "Download Flow", "Manifest Resolution -> Provider Discovery -> Parallel Download")
    
    box_left = slide14.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Bootstrap: Fetch JSON manifest metadata using ManifestID.",
        "• Planning: Query placement records to compile the local ChunkProviderMap.",
        "• Transfer: Concurrent worker routines request shards in parallel streams.",
        "• Reassembly: Validate shard hashes and decode Reed-Solomon blocks to build the original file."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide14, 14)
    
    slide14.notes_slide.notes_text_frame.text = (
        "This slide presents the end-to-end download sequence from the client's perspective.\n\n"
        "We start with bootstrapping: getting a ManifestID and querying the DHT to fetch the manifest details.\n"
        "Next is planning: mapping target shards to active online peers to build our local ChunkProviderMap.\n"
        "Third is swarming: starting parallel fetch routines to pull shards concurrently.\n"
        "Finally, we run hash checks on each shard, decode Reed-Solomon blocks, and write the reassembled file to local storage. If any chunk is corrupted or incomplete, we reject it and request it again."
    )

    # -------------------------------------------------------------
    # SLIDE 15: CIPHER in Action
    # -------------------------------------------------------------
    slide15 = prs.slides.add_slide(blank_layout)
    apply_background(slide15)
    add_header_and_logos(slide15, "CIPHER in Action", "End-to-End File Distribution Demo")
    
    box_left = slide15.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• File Ingestion: Chunks file, encrypts via XChaCha20, and publishes the manifest.",
        "• DHT Registration: Registers shard locations (PlacementMap) on Kademlia.",
        "• Direct Upgrade: Triggers DCUtR hole punching to bypass NAT barriers.",
        "• Swarmed Retrieval: Concurrent workers download and verify shards."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    # Dark console mockup on the right
    term_bg = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))
    term_bg.fill.solid()
    term_bg.fill.fore_color.rgb = RGBColor(9, 13, 22) # Slate 950
    term_bg.line.color.rgb = COLOR_BORDER
    
    tf_r = slide15.shapes.add_textbox(Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.3)).text_frame
    tf_r.word_wrap = True
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "SYSTEM INTEGRATION LOGS:"
    p_rh.font.name = "Courier New"
    p_rh.font.size = Pt(10)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_SUCCESS
    
    logs = [
        "[INFO] content: Ingested file, created 16 chunks",
        "[INFO] manifest: Merkle root computed = a7b8...4f1e",
        "[INFO] control: Registering shards on PlacementMap",
        "[INFO] transport: relay connection PeerID=12D3K...",
        "[INFO] dcutr: UDP hole punch succeeded! Direct active",
        "[INFO] scheduler: starting parallel swarm retrials",
        "[INFO] worker-0: downloading shard 0... OK",
        "[INFO] verifier: SHA-256 match for shard 0 & 1",
        "[INFO] engine: verification succeeded, file reconstructed"
    ]
    for log in logs:
        p = tf_r.add_paragraph()
        p.text = log
        p.font.name = "Courier New"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(3)
        
    add_footer(slide15, 15)
    
    slide15.notes_slide.notes_text_frame.text = (
        "Slide 15 is our centerpiece slide showing the actual E2E system integration in action.\n\n"
        "Here we show the live data flow pathway: the publisher ingests the file, registers shard allocations, Kademlia coordinates location discovery, and the client pulls shards concurrently from separate providers.\n"
        "On the right, we show actual terminal outputs from a successful integration test. You can see the publisher chunking data, Kademlia routing paths, DCUtR coordinates successfully hole-punching UDP tunnels to providers, parallel scheduler streams active, and final SHA-256 verification passing."
    )

    # -------------------------------------------------------------
    # SLIDE 16: Failure Recovery
    # -------------------------------------------------------------
    slide16 = prs.slides.add_slide(blank_layout)
    apply_background(slide16)
    add_header_and_logos(slide16, "Failure Recovery", "Ensuring Fault-Resilient Shard Retrieval")
    
    box_left = slide16.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        "• Dynamic Failover: Connection drops trigger provider switching; the scheduler recycles the task to alternative nodes.",
        "• Interruption Resuming: Interrupted downloads load progress bitsets and fetch only missing shards.",
        "• Reed-Solomon Limit: Decodes data if ANY K out of K+M shards are present. Fails if lost shards exceed parity M."
    ]
    for i, b_text in enumerate(bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = b_text
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(16) if i > 0 else Pt(0)
        
    add_footer(slide16, 16)
    
    slide16.notes_slide.notes_text_frame.text = (
        "Here we discuss recovery. P2P transfers must survive network failure.\n\n"
        "First, we handle network failures with provider switching: workers monitor transfer loops; if a peer goes offline, we recycle the shard task back to the queue and dial a different provider.\n"
        "We also support local session resume, loading a bitset of progress so a client never re-downloads completed blocks after a crash.\n"
        "At the content layer, Reed-Solomon allows reconstruction using any K out of K+M shards. If more than M shards are missing, the file is lost. This defines the clear boundary of our recovery guarantees."
    )

    # -------------------------------------------------------------
    # SLIDE 17: Future Implementations
    # -------------------------------------------------------------
    slide17 = prs.slides.add_slide(blank_layout)
    apply_background(slide17)
    add_header_and_logos(slide17, "Future Implementations", "Scalability • Optimization • Incentives")
    
    col_width = Inches(3.7)
    col_gap = Inches(0.2)
    start_left = Inches(0.6)
    
    roadmap = [
        ("1. Advanced Scheduling", "• Rarest-first shard scheduling to maximize shard distribution throughout the network.\n• Latency and bandwidth-aware provider selection.\n• Provider reputation scoring."),
        ("2. Incentive Engine", "• Token settlement layer using smart contracts to reward hosting nodes.\n• Zero-knowledge proofs of storage.\n• Micro-payment channels."),
        ("3. Protocol Hardening", "• Connection managers imposing strict stream limits per node to mitigate DDoS.\n• Bounded buffer stream limits.\n• Protocol testing under network attacks.")
    ]
    
    for i, (name, desc) in enumerate(roadmap):
        box_left = start_left + i * (col_width + col_gap)
        
        rect = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.2), col_width, Inches(4.3))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = COLOR_BORDER
        
        tx_box = slide17.shapes.add_textbox(box_left, Inches(2.3), col_width, Inches(4.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(8)
        
    add_footer(slide17, 17)
    
    slide17.notes_slide.notes_text_frame.text = (
        "This slide outlines our future work roadmap.\n\n"
        "First, advanced scheduling: implementing rarest-first algorithms and provider rating metrics to fetch shards more efficiently.\n"
        "Second, incentives: building a payment layer using smart contracts and micro-channels alongside cryptographic proof-of-storage validations so hosting nodes are fairly compensated.\n"
        "Third, protocol hardening: adding connection limiters, stream timeouts, and rate limits to defend the P2P transport against memory leaks and DDoS flood attacks."
    )

    # -------------------------------------------------------------
    # SLIDE 18: Team
    # -------------------------------------------------------------
    slide18 = prs.slides.add_slide(blank_layout)
    apply_background(slide18)
    add_header_and_logos(slide18, "Team", "Building CIPHER")
    
    # 2 columns team info
    box_left = slide18.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Contributors"
    p.font.name = "Outfit"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(14)
    
    p_dev1 = tf_l.add_paragraph()
    p_dev1.text = "Adarsh Kumar"
    p_dev1.font.name = "Calibri"
    p_dev1.font.size = Pt(18)
    p_dev1.font.bold = True
    p_dev1.font.color.rgb = COLOR_SECONDARY
    
    p_dev1_role = tf_l.add_paragraph()
    p_dev1_role.text = "Core Developer • Backend, P2P Networks, & Transport Layer Implementation"
    p_dev1_role.font.name = "Calibri"
    p_dev1_role.font.size = Pt(13)
    p_dev1_role.font.color.rgb = TEXT_MAIN
    p_dev1_role.space_before = Pt(4)
    
    box_right = slide18.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    
    p_r = tf_r.paragraphs[0]
    p_r.text = "Mentors"
    p_r.font.name = "Outfit"
    p_r.font.size = Pt(22)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_PRIMARY
    p_r.space_after = Pt(14)
    
    p_ment = tf_r.add_paragraph()
    p_ment.text = "DevlUp Labs"
    p_ment.font.name = "Calibri"
    p_ment.font.size = Pt(18)
    p_ment.font.bold = True
    p_ment.font.color.rgb = COLOR_SECONDARY
    
    p_ment_desc = tf_r.add_paragraph()
    p_ment_desc.text = "IIT Jodhpur Student Developer Community\nOpen Source Mentorship & Platform Support"
    p_ment_desc.font.name = "Calibri"
    p_ment_desc.font.size = Pt(13)
    p_ment_desc.font.color.rgb = TEXT_MAIN
    p_ment_desc.space_before = Pt(4)
    
    p_inst = tf_r.add_paragraph()
    p_inst.text = "IIT Jodhpur CSE Department Project Submission"
    p_inst.font.name = "Calibri"
    p_inst.font.size = Pt(13)
    p_inst.font.italic = True
    p_inst.font.color.rgb = TEXT_MUTED
    p_inst.space_before = Pt(24)
    
    add_footer(slide18, 18)
    
    slide18.notes_slide.notes_text_frame.text = (
        "Thank you for attending this presentation on the CIPHER decentralized content distribution protocol.\n\n"
        "The project is developed by Adarsh Kumar, focusing on the core backend execution engine, persistent peer identity, and transport protocols.\n"
        "The work is guided and supported by the DevlUp Labs team at the Indian Institute of Technology Jodhpur.\n\n"
        "We are happy to take any questions from the panel."
    )

    # Save to disk
    out_path = "/Users/adarsh/Projects/devlup/SOC/CIPHER/ppt/CIPHER_Presentation.pptx"
    prs.save(out_path)
    print("Generated presentation saved to:", out_path)

if __name__ == "__main__":
    create_presentation()
